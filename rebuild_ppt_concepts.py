# -*- coding: utf-8 -*-
"""
[이 스크립트가 하는 일]
기존 PPT 파일(01_stock_prediction_report.pptx)에
머신러닝 모델 개념 소개 슬라이드 6장을 자동으로 삽입합니다.

삽입 위치: XGBoost 슬라이드(9번) 바로 뒤 → 슬라이드 10~15번 위치
삽입 내용: 선형 회귀 / Ridge / Lasso / ElasticNet / 결정 트리 / KNN

실행 방법:
    python rebuild_ppt_concepts.py
결과물:
    01_stock_prediction_report.pptx (덮어쓰기)
"""

# ─────────────────────────────────────────────────────────────
# 라이브러리 불러오기
# ─────────────────────────────────────────────────────────────
from pptx import Presentation          # PPT 파일 읽기/쓰기
from pptx.util import Inches, Pt       # 크기 단위 변환 (인치, 포인트)
from pptx.dml.color import RGBColor    # RGB 색상 지정
from pptx.enum.text import PP_ALIGN    # 텍스트 정렬 방식 (왼쪽/가운데/오른쪽)

# ─────────────────────────────────────────────────────────────
# XML 네임스페이스 상수
# ─────────────────────────────────────────────────────────────
# PPT 파일(.pptx)은 내부적으로 XML로 구성되어 있습니다.
# 슬라이드 순서를 바꾸거나 삭제할 때 XML 요소를 직접 조작해야 하는데,
# 그때 각 XML 태그가 어느 '규격(네임스페이스)'에 속하는지 명시해야 합니다.
NS_RELS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
# → 슬라이드 간 관계(링크) 정보가 담긴 XML 네임스페이스

NS_PML  = 'http://schemas.openxmlformats.org/presentationml/2006/main'
# → 슬라이드 목록(순서) 정보가 담긴 XML 네임스페이스


# ─────────────────────────────────────────────────────────────
# PPT 파일 열기
# ─────────────────────────────────────────────────────────────
prs = Presentation('01_stock_prediction_report.pptx')

# 나중에 슬라이드 순서를 재배치할 때 "원래 슬라이드 수"가 필요하므로 미리 저장
n_original = len(prs.slides)
print('로드 시 슬라이드 수:', n_original)

# 슬라이드 순서 목록을 담고 있는 XML 요소를 변수에 담아둡니다.
# (나중에 순서 재배치할 때 이 객체를 직접 조작합니다)
sldIdLst = prs.slides._sldIdLst


# ─────────────────────────────────────────────────────────────
# 공통 헬퍼 함수 정의
# ─────────────────────────────────────────────────────────────

def rgb(r, g, b):
    """R, G, B 숫자 세 개를 받아 python-pptx가 인식하는 색상 객체로 변환."""
    return RGBColor(r, g, b)


# ── 슬라이드 전체에서 사용하는 색상 팔레트 ────────────────────
COLOR_BG    = rgb(15, 23, 42)    # 슬라이드 배경색 (어두운 네이비)
COLOR_WHITE = rgb(255, 255, 255) # 일반 텍스트
COLOR_GRAY  = rgb(148, 163, 184) # 하단 푸터 텍스트
COLOR_GREEN = rgb(74, 222, 128)  # 장점 패널 제목
COLOR_RED   = rgb(248, 113, 113) # 단점 패널 제목

# 모델 계열별 포인트 컬러 (헤더 바, 강조 텍스트에 사용)
CAT_COLORS = {
    'linear': rgb(234, 168, 23),   # 선형 계열 → 금색
    'tree':   rgb(167, 139, 250),  # 트리 계열 → 보라색
    'knn':    rgb(34, 211, 238),   # 거리 기반  → 하늘색
}


def add_rect(slide, left, top, width, height, fill):
    """
    슬라이드에 색이 채워진 직사각형을 추가합니다.

    매개변수:
        slide  : 도형을 추가할 슬라이드 객체
        left   : 왼쪽 가장자리에서의 거리 (인치)
        top    : 위쪽 가장자리에서의 거리 (인치)
        width  : 직사각형의 가로 길이 (인치)
        height : 직사각형의 세로 길이 (인치)
        fill   : 채울 색상 (RGBColor 객체)
    반환값:
        생성된 도형(shape) 객체
    """
    # 도형 타입 1 = 직사각형(MSO_SHAPE_TYPE.RECTANGLE)
    s = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid()                 # 단색으로 채우기
    s.fill.fore_color.rgb = fill   # 채울 색상 지정
    s.line.fill.background()       # 테두리 선 제거 (배경과 동화)
    return s


def add_tb(slide, left, top, width, height,
           text, size=12, bold=False, color=None,
           align=PP_ALIGN.LEFT, wrap=True):
    """
    슬라이드에 텍스트 박스를 추가합니다.

    매개변수:
        slide  : 텍스트 박스를 추가할 슬라이드 객체
        left   : 왼쪽 가장자리에서의 거리 (인치)
        top    : 위쪽 가장자리에서의 거리 (인치)
        width  : 텍스트 박스 가로 길이 (인치)
        height : 텍스트 박스 세로 길이 (인치)
        text   : 표시할 문자열
        size   : 폰트 크기 (pt 단위, 기본값 12)
        bold   : 굵게 여부 (기본값 False)
        color  : 글자 색상 (None이면 흰색 적용)
        align  : 텍스트 정렬 (기본값 왼쪽)
        wrap   : 텍스트 자동 줄바꿈 여부 (기본값 True)
    반환값:
        생성된 텍스트 박스 객체
    """
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap            # 박스 너비를 넘으면 자동 줄바꿈
    p = tf.paragraphs[0]           # 텍스트 박스의 첫 번째 단락
    p.alignment = align
    r = p.add_run()                # 단락 안에 텍스트 런(run) 추가
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or COLOR_WHITE  # 색상 미지정 시 흰색
    return tb


def set_bg(slide, color):
    """슬라이드 전체 배경색을 단색으로 설정합니다."""
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def make_slide(prs, title, sub, tagline,
               an_title, an_body,
               hw_title, hw_items,
               pros, cons,
               page, accent, badge=''):
    """
    모델 개념 소개 슬라이드 1장을 생성해 prs에 추가합니다.
    (add_slide는 항상 맨 뒤에 붙습니다. 순서 조정은 STEP 3에서 합니다.)

    슬라이드 레이아웃 (가로 10인치 × 세로 7인치 기준):
    ┌──────────────────────────────────────────┐  ← 0~1.1인치
    │ [배지]  모델 이름 (title)           페이지 │  ← 헤더 바 (accent 색상)
    │          부제목 (sub)                      │
    ├──────────────────────────────────────────┤  ← 1.1~1.62인치
    │  한 줄 요약  |  tagline                   │  ← 요약 바 (어두운 색)
    ├─────────────────┬────────────────────────┤  ← 1.75~4.65인치
    │ 비유로 이해하기  │  어떻게 작동하나        │  ← 좌/우 패널
    │  (an_body)      │  (hw_items 목록)        │
    ├─────────────────┼────────────────────────┤  ← 4.78~6.6인치
    │ ✅ 장점 (pros)  │  ⚠ 단점 (cons)         │  ← 장단점 패널
    ├──────────────────────────────────────────┤  ← 6.78인치~
    │ 프로젝트 이름                      페이지 │  ← 하단 푸터
    └──────────────────────────────────────────┘

    매개변수:
        prs      : Presentation 객체 (슬라이드가 여기에 추가됨)
        title    : 모델 이름 (예: '선형 회귀 (Linear Regression)')
        sub      : 한 줄 부제목
        tagline  : 요약 바에 표시할 핵심 메시지
        an_title : 비유 패널 제목 (보통 '비유로 이해하기')
        an_body  : 비유 패널 본문 (여러 줄 문자열)
        hw_title : 작동방식 패널 제목 (보통 '어떻게 작동하나')
        hw_items : 작동방식 항목 리스트 (각 항목이 한 줄씩 출력됨)
        pros     : 장점 리스트
        cons     : 단점 리스트
        page     : 슬라이드 번호 표시 문자열 (예: '10 / 22')
        accent   : 이 슬라이드의 포인트 컬러 (CAT_COLORS 값 사용)
        badge    : 헤더 좌측 배지 텍스트 (예: '선형 계열')
    """
    # 슬라이드 레이아웃 0번 사용 (이 PPT에 레이아웃이 1개뿐이라 고정)
    layout = prs.slide_layouts[0]
    slide  = prs.slides.add_slide(layout)

    # 레이아웃에서 자동으로 생긴 제목/내용 placeholder를 모두 제거합니다.
    # → 이 스크립트는 모든 요소를 직접 배치하므로 기본 틀이 방해가 됩니다.
    for ph in slide.placeholders:
        sp = ph._element
        sp.getparent().remove(sp)

    # 배경색 설정
    set_bg(slide, COLOR_BG)

    # ── 헤더 바 (슬라이드 상단, accent 컬러 직사각형) ────────────
    add_rect(slide, 0, 0, 10, 1.1, accent)

    # 배지 (좌측 상단의 계열 표시 박스, 예: "선형 계열")
    if badge:
        # 배지 배경: 슬라이드 배경색으로 작은 박스 → 포인트 컬러 글자와 대비
        add_rect(slide, 0.18, 0.16, 1.3, 0.42, rgb(15, 23, 42))
        add_tb(slide, 0.18, 0.16, 1.3, 0.42,
               badge, size=9, bold=True, color=accent,
               align=PP_ALIGN.CENTER)

    # 모델 이름 (큰 글씨, 헤더 바 위에 어두운 색으로)
    add_tb(slide, 1.62, 0.08, 7.0, 0.58,
           title, size=21, bold=True,
           color=rgb(15, 23, 42))

    # 부제목 (모델 이름 아래 작은 글씨)
    add_tb(slide, 1.62, 0.64, 7.0, 0.36,
           sub, size=10, color=rgb(30, 41, 59))

    # 페이지 번호 (우측 상단)
    add_tb(slide, 8.7, 0.08, 1.1, 0.5,
           page, size=10, color=rgb(30, 41, 59),
           align=PP_ALIGN.RIGHT)

    # ── 한 줄 요약 바 (헤더 바 아래 가로 전체) ───────────────────
    add_rect(slide, 0, 1.1, 10, 0.52, rgb(30, 41, 59))
    add_tb(slide, 0.25, 1.15, 9.5, 0.42,
           f'한 줄 요약   |   {tagline}',
           size=11, color=accent)

    # ── 왼쪽 패널: 비유로 이해하기 ───────────────────────────────
    add_rect(slide, 0.18, 1.75, 4.65, 2.9, rgb(22, 33, 55))  # 패널 배경
    add_tb(slide, 0.32, 1.82, 4.35, 0.38,
           an_title, size=11, bold=True, color=accent)         # 패널 제목
    add_tb(slide, 0.32, 2.24, 4.35, 2.32,
           an_body, size=10.5, color=COLOR_WHITE)              # 비유 본문

    # ── 오른쪽 패널: 어떻게 작동하나 ────────────────────────────
    add_rect(slide, 5.17, 1.75, 4.65, 2.9, rgb(22, 33, 55))  # 패널 배경
    add_tb(slide, 5.32, 1.82, 4.35, 0.38,
           hw_title, size=11, bold=True, color=accent)         # 패널 제목

    # 작동 단계 항목을 세로로 나열 (각 항목마다 y 위치를 0.44인치씩 내림)
    y = 2.24
    for item in hw_items:
        add_tb(slide, 5.32, y, 4.35, 0.4,
               item, size=10, color=COLOR_WHITE)
        y += 0.44

    # ── 왼쪽 하단: 장점 패널 ─────────────────────────────────────
    add_rect(slide, 0.18, 4.78, 4.65, 1.82, rgb(16, 38, 28))  # 녹색 계열 배경
    add_tb(slide, 0.32, 4.85, 2.0, 0.36,
           '✅  장점', size=11, bold=True, color=COLOR_GREEN)

    # 장점 항목을 세로로 나열 (각 항목마다 y 위치를 0.4인치씩 내림)
    y = 5.24
    for p_text in pros:
        add_tb(slide, 0.32, y, 4.4, 0.38,
               f'• {p_text}', size=10, color=COLOR_WHITE)
        y += 0.4

    # ── 오른쪽 하단: 단점 패널 ───────────────────────────────────
    add_rect(slide, 5.17, 4.78, 4.65, 1.82, rgb(38, 16, 16))  # 붉은 계열 배경
    add_tb(slide, 5.32, 4.85, 2.0, 0.36,
           '⚠  단점', size=11, bold=True, color=COLOR_RED)

    # 단점 항목을 세로로 나열
    y = 5.24
    for c_text in cons:
        add_tb(slide, 5.32, y, 4.4, 0.38,
               f'• {c_text}', size=10, color=COLOR_WHITE)
        y += 0.4

    # ── 하단 푸터 바 ─────────────────────────────────────────────
    add_rect(slide, 0, 6.78, 10, 0.22, rgb(22, 33, 55))
    add_tb(slide, 0.2, 6.8, 7.5, 0.18,
           '주식 등락률 예측 AI 모델 개발', size=8, color=COLOR_GRAY)
    add_tb(slide, 8.0, 6.8, 1.8, 0.18,
           page, size=8, color=COLOR_GRAY, align=PP_ALIGN.RIGHT)

    return slide


# ─────────────────────────────────────────────────────────────
# STEP 1: 모델 개념 슬라이드 6장 생성
#
# python-pptx의 add_slide()는 항상 맨 뒤에 슬라이드를 붙입니다.
# 그래서 일단 6장을 모두 맨 뒤에 추가한 뒤,
# STEP 2에서 원하는 위치(XGBoost 슬라이드 뒤)로 순서를 재배치합니다.
# ─────────────────────────────────────────────────────────────

# ── 1. 선형 회귀 ─────────────────────────────────────────────
make_slide(
    prs,
    title    = '선형 회귀 (Linear Regression)',
    sub      = '가장 단순한 예측 공식 — 다른 모델 성능의 기준선(Baseline) 역할',
    tagline  = '"각 지표에 가중치를 곱해 더하면 등락률이 나온다"',
    an_title = '비유로 이해하기',
    an_body  = (
        '키와 몸무게의 관계처럼,\n'
        '"RSI가 1 오르면 등락률 0.03% 변화"\n'
        '이런 공식을 데이터에서 자동으로 찾습니다.\n\n'
        '수만 개의 과거 날짜를 보면서\n'
        '"어떤 가중치 조합이 정답에 가장\n'
        ' 가까운지"를 계산합니다.\n\n'
        '예:  y = 0.03×RSI + 0.02×MACD\n'
        '       + 0.01×거래량 + ...'
    ),
    hw_title = '어떻게 작동하나',
    hw_items = [
        '① 11개 지표 각각에 가중치를 부여',
        '② 가중치×지표를 모두 더함 = 예측값',
        '③ 오차를 최소화하는 가중치를 학습',
        '④ 가중치를 보면 어떤 지표가',
        '   중요한지 바로 알 수 있음',
    ],
    pros     = [
        '학습 시간 0.1초 미만 — 9개 중 가장 빠름',
        '가중치를 보면 지표별 영향력 즉시 파악',
        '다른 모델의 성능 비교 기준선 역할',
    ],
    cons     = [
        '"RSI 높고 거래량 급증" 같은 조합 효과 미반영',
        '비선형 복잡한 시장 패턴 표현 한계',
        '규제 없으면 과적합 위험 (→ Ridge/Lasso)',
    ],
    page     = '10 / 18',
    accent   = CAT_COLORS['linear'],
    badge    = '선형 계열',
)

# ── 2. Ridge 회귀 ────────────────────────────────────────────
make_slide(
    prs,
    title    = 'Ridge 회귀',
    sub      = '선형 회귀 + "너무 극단적인 답을 내지 마" 제약',
    tagline  = '"모든 가중치를 조금씩 줄여 안정적인 예측을 만든다"',
    an_title = '비유로 이해하기',
    an_body  = (
        '"답안을 너무 극단적으로 쓰지 마세요"\n'
        '라는 채점 규정을 상상해 보세요.\n\n'
        '선형 회귀는 가중치가 지나치게 클 수 있어요.\n'
        '→ "이 지표가 무조건 1,000% 중요해!"\n\n'
        'Ridge는 가중치가 커질수록 페널티를 주어\n'
        '"모든 지표를 균형 있게 활용"하도록 합니다.\n\n'
        'Alpha(α)가 클수록 가중치를 더 강하게 억제\n'
        '→ 탐색 결과: α = 0.1 선택'
    ),
    hw_title = '어떻게 작동하나',
    hw_items = [
        '① 선형 회귀처럼 가중치 학습',
        '② 가중치² 합산에 페널티 추가',
        '③ Alpha(α)로 페널티 강도 조절',
        '④ 모든 가중치를 0이 아닌',
        '   작은 값으로 유지 (완전 제거 X)',
    ],
    pros     = [
        '과적합 방지 — 선형 회귀보다 안정적',
        '모든 지표를 균형 있게 조금씩 활용',
        '학습 시간 1초 미만',
    ],
    cons     = [
        '불필요한 지표도 완전히 제거하지 않음',
        'Alpha 최적값을 별도 탐색해야 함',
        '비선형 관계 표현 불가',
    ],
    page     = '11 / 18',
    accent   = CAT_COLORS['linear'],
    badge    = '선형 계열',
)

# ── 3. Lasso 회귀 ────────────────────────────────────────────
make_slide(
    prs,
    title    = 'Lasso 회귀',
    sub      = '중요하지 않은 지표를 완전히 0으로 만들어 자동 정리',
    tagline  = '"불필요한 지표를 자동으로 걸러내 진짜 중요한 것만 남긴다"',
    an_title = '비유로 이해하기',
    an_body  = (
        '11명이 투자 의견을 말하는데,\n'
        '"설득력 없는 의견은 아예 발언권 0으로"\n'
        '만드는 사회자를 상상해 보세요.\n\n'
        'Ridge: 11명 모두 발언, 극단적 주장만 줄임\n'
        'Lasso: 핵심 멤버만 남기고 나머지 침묵\n\n'
        '→ Lasso 결과: 예측에 중요한 지표만 남고\n'
        '   나머지 가중치가 정확히 0이 됩니다.\n\n'
        '→ 탐색 결과: α = 0.001 선택'
    ),
    hw_title = '어떻게 작동하나',
    hw_items = [
        '① 선형 회귀처럼 가중치 학습',
        '② 가중치 절댓값 합산에 페널티 추가',
        '③ 페널티가 강해지면 일부 가중치 → 0',
        '④ 가중치 0인 지표 = 예측에 불필요',
        '   → 자동으로 변수 선택',
    ],
    pros     = [
        '자동 변수 선택 — 중요한 지표가 무엇인지 밝힘',
        '불필요한 지표를 완전 제거해 모델 단순화',
        '과적합 방지 + 해석 가능성 동시 확보',
    ],
    cons     = [
        '상관관계 높은 지표 중 하나를 임의 제거 가능',
        'Alpha가 너무 크면 모든 가중치가 0이 됨',
        '비선형 관계 표현 불가',
    ],
    page     = '12 / 18',
    accent   = CAT_COLORS['linear'],
    badge    = '선형 계열',
)

# ── 4. ElasticNet 회귀 ───────────────────────────────────────
make_slide(
    prs,
    title    = 'ElasticNet 회귀',
    sub      = 'Ridge + Lasso를 섞은 하이브리드 모델',
    tagline  = '"Ridge의 안정성 + Lasso의 변수 선택을 동시에 얻는다"',
    an_title = '비유로 이해하기',
    an_body  = (
        '두 가지 다이어트 방법이 있다고 해요.\n\n'
        'Ridge식: 모든 음식을 조금씩 줄이기\n'
        'Lasso식: 나쁜 음식은 완전히 끊기\n\n'
        'ElasticNet은 "둘을 반반씩" 하는 방식.\n\n'
        'l1_ratio = 0.5\n'
        '→ Ridge 50% + Lasso 50% 동시 적용\n\n'
        '→ 탐색 결과: α=0.001, l1_ratio=0.5 선택'
    ),
    hw_title = '어떻게 작동하나',
    hw_items = [
        '① Ridge 페널티 + Lasso 페널티 동시 적용',
        '② l1_ratio로 두 방식의 비율 조절',
        '③ 일부 계수 → 0 (Lasso 효과)',
        '④ 나머지 계수는 고르게 억제 (Ridge 효과)',
        '⑤ 두 하이퍼파라미터(α, l1_ratio) 탐색',
    ],
    pros     = [
        'Ridge + Lasso 장점 동시에 누림',
        '지표 간 상관관계가 높을 때 특히 효과적',
        '가장 유연한 선형 모델',
    ],
    cons     = [
        '탐색해야 할 하이퍼파라미터가 2개',
        'Ridge나 Lasso보다 해석이 약간 복잡',
        '비선형 관계 표현 불가',
    ],
    page     = '13 / 18',
    accent   = CAT_COLORS['linear'],
    badge    = '선형 계열',
)

# ── 5. 결정 트리 ─────────────────────────────────────────────
make_slide(
    prs,
    title    = '결정 트리 (Decision Tree)',
    sub      = '조건 분기를 반복해 예측하는 규칙 기반 모델',
    tagline  = '"RSI가 낮고 거래량이 높으면 상승" — 사람이 읽을 수 있는 규칙',
    an_title = '비유로 이해하기',
    an_body  = (
        '병원 응급실 환자 분류를 떠올려 보세요.\n\n'
        '"체온 > 38도?"\n'
        '  Yes → "혈압 정상?"\n'
        '          No → "즉시 입원"\n\n'
        '주식도 마찬가지:\n'
        '"RSI < 30?" (과매도)\n'
        '  Yes → "거래량 급증?"\n'
        '          Yes → "상승 +2.3% 예측"\n\n'
        'max_depth=8: 분기를 8번으로 제한\n'
        '→ 너무 깊어지면 훈련 데이터만 외움'
    ),
    hw_title = '어떻게 작동하나',
    hw_items = [
        '① 어떤 지표를 어떤 기준으로 나눌지 학습',
        '② 나눌 때마다 오차가 가장 줄어드는 분기',
        '③ max_depth=8 에서 분기 중단',
        '④ 마지막 그룹(Leaf)의 평균 = 예측값',
        '⑤ 규칙을 출력하면 사람이 직접 읽기 가능',
    ],
    pros     = [
        '예측 규칙을 직접 읽을 수 있어 해석성 최고',
        '데이터 전처리(정규화 등) 불필요',
        '비선형 관계 및 지표 간 조합 포착 가능',
    ],
    cons     = [
        '깊어질수록 훈련 데이터를 외워버림 (과적합)',
        '9개 모델 중 2025년 실전 과적합 갭 가장 큼 (+0.7%p)',
        'RF/XGB보다 성능 낮음 (트리 1개의 한계)',
    ],
    page     = '14 / 18',
    accent   = CAT_COLORS['tree'],
    badge    = '트리 계열',
)

# ── 6. KNN 회귀 ──────────────────────────────────────────────
make_slide(
    prs,
    title    = 'KNN 회귀 (K-Nearest Neighbors)',
    sub      = '"오늘과 가장 비슷했던 과거 날들의 결과 평균" = 예측',
    tagline  = '"오늘 날씨 같은 날 다음엔 어땠지?" — 유사한 과거 k개를 참조',
    an_title = '비유로 이해하기',
    an_body  = (
        '동네 부동산 가격 예측을 생각해 보세요.\n\n'
        '"이 집과 비슷한 집 20채를 찾아보자.\n'
        ' 그 집들의 실거래가 평균이 이 집 가격."\n\n'
        '주식도 마찬가지:\n'
        '오늘: RSI=42, 거래량=1.2배, MACD=0.3...\n'
        '→ 과거에 이런 상황이었던 날 20개 검색\n'
        '→ 그 날들의 5일 후 등락률 평균 = 예측\n\n'
        'k=20: 가장 비슷한 20개 날 참조\n'
        '(탐색에서 k=20일 때 MAE 최소)'
    ),
    hw_title = '어떻게 작동하나',
    hw_items = [
        '① 학습 = 훈련 데이터를 통째로 저장',
        '② 예측 = 오늘과 거리 가장 가까운',
        '   k=20개 날 탐색 (11개 지표 기준)',
        '③ 20개 날의 5일 후 등락률 평균 = 예측',
        '④ 데이터 많을수록 참조 후보 증가',
    ],
    pros     = [
        '학습 과정 없음 — 즉시 사용 가능',
        '"비슷한 과거 = 비슷한 미래" 직관적 원리',
        '데이터 추가만 하면 자동으로 모델 개선',
    ],
    cons     = [
        '예측마다 전체 훈련 데이터 탐색 → 느림',
        '11개 지표가 모두 동등하게 취급됨',
        '지표 수가 늘면 거리 의미 희석 (차원의 저주)',
    ],
    page     = '15 / 18',
    accent   = CAT_COLORS['knn'],
    badge    = '거리 기반',
)

print('슬라이드 생성 후 총 수:', len(prs.slides))


# ─────────────────────────────────────────────────────────────
# STEP 2: 새 슬라이드 6장을 슬라이드 9번(XGBoost) 바로 뒤로 이동
#
# [문제 상황]
# add_slide()는 항상 맨 뒤에 붙습니다.
# 지금 상태: [기존 0~8] [기존 9~끝] [새 슬라이드 6장]
#
# [목표 상태]
# [기존 0~8] [새 슬라이드 6장] [기존 9~끝]
#
# [방법]
# PPT의 슬라이드 순서는 XML 내 <sldIdLst> 태그의 자식 요소 순서로 결정됩니다.
# 이 자식 요소들을 원하는 순서로 재배치하면 슬라이드 순서가 바뀝니다.
# ─────────────────────────────────────────────────────────────

# 현재 XML에 있는 모든 슬라이드 ID 요소를 순서대로 가져옵니다.
all_ids = sldIdLst.findall(f'{{{NS_PML}}}sldId')

# 원하는 새 순서로 분류합니다.
old_first9 = all_ids[:9]            # 기존 슬라이드 앞 9장 (슬라이드 1~9번)
new_six    = all_ids[n_original:]   # 방금 추가된 개념 슬라이드 6장
old_rest   = all_ids[9:n_original]  # 기존 슬라이드 나머지 (슬라이드 10번~끝)

# 기존 순서를 모두 제거한 뒤, 원하는 순서로 다시 추가합니다.
for elem in all_ids:
    sldIdLst.remove(elem)
for elem in old_first9 + new_six + old_rest:
    sldIdLst.append(elem)

print('순서 조정 후 총 슬라이드:', len(prs.slides))


# ─────────────────────────────────────────────────────────────
# STEP 3: 파일 저장 (기존 파일 덮어쓰기)
# ─────────────────────────────────────────────────────────────
prs.save('01_stock_prediction_report.pptx')
print('PPT saved OK')
